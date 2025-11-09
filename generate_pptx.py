#!/usr/bin/env python3
"""
Script to generate a PowerPoint presentation from PNG images in the slides directory.
Each PNG becomes its own slide with the image sized to fit the slide.
"""

import os
import glob
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

def create_presentation_from_images(slides_dir='slides', output_file='adc_talk.pptx'):
    """Create a PowerPoint presentation with each PNG image as a separate slide."""

    # Get all PNG files and sort them alphabetically
    png_files = sorted(glob.glob(os.path.join(slides_dir, '*.png')))

    if not png_files:
        print(f"No PNG files found in {slides_dir}/")
        return

    print(f"Found {len(png_files)} PNG files")

    # Create presentation with 16:9 aspect ratio (standard for modern presentations)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Process each PNG file
    for idx, png_file in enumerate(png_files, 1):
        filename = os.path.basename(png_file)
        print(f"Processing {idx}/{len(png_files)}: {filename}")

        # Add a blank slide
        blank_slide_layout = prs.slide_layouts[6]  # 6 = blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Get image dimensions to calculate proper sizing
        with Image.open(png_file) as img:
            img_width, img_height = img.size
            img_aspect_ratio = img_width / img_height

        # Calculate dimensions to fit image within slide while maintaining aspect ratio
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        slide_aspect_ratio = slide_width / slide_height

        if img_aspect_ratio > slide_aspect_ratio:
            # Image is wider than slide - fit to width
            pic_width = slide_width
            pic_height = int(slide_width / img_aspect_ratio)
            left = 0
            top = (slide_height - pic_height) // 2
        else:
            # Image is taller than slide - fit to height
            pic_height = slide_height
            pic_width = int(slide_height * img_aspect_ratio)
            left = (slide_width - pic_width) // 2
            top = 0

        # Add picture to slide
        slide.shapes.add_picture(png_file, left, top, width=pic_width, height=pic_height)

    # Save presentation
    prs.save(output_file)
    print(f"\nPresentation saved as: {output_file}")
    print(f"Total slides: {len(png_files)}")

if __name__ == '__main__':
    create_presentation_from_images()
